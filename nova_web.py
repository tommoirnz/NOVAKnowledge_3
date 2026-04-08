"""
Nova Web Interface -
 NOVA — Full Feature Set
  * Web TTS (edge_tts) + mute button
  * Voice input (Whisper)
  * MathJax rendering
  * marked.js markdown
  * Code-block copy buttons
  * Inline image / diagram / audio / video
  * Save-response button (download or new tab)
  * Drag-and-drop file upload (full preview)
  * Working online/offline status dot
"""

import http.server
import socketserver
import json
import threading
import webbrowser
import os
from urllib.parse import urlparse
from datetime import datetime

_history_lock = threading.Lock()


# =============================================================================
#  HTTP HANDLER
# =============================================================================

class NovaWebHandler(http.server.SimpleHTTPRequestHandler):
    nova = None

    # ── GET ──────────────────────────────────────────────────────────────────
    def do_GET(self):
        path = urlparse(self.path).path

        if path == '/':
            self._serve_html()
        elif path == '/api/history':
            self._serve_json(self.nova.conversation_history)
        elif path == '/api/state':
            self._serve_json({
                'model':    self.nova.ai.model,
                'thinking': getattr(self.nova, '_thinking', False),
                'tts':      getattr(self.nova, '_tts_on', False),
                'web_tts':  getattr(self.nova, '_web_tts_on', True),
            })
        elif path == '/api/ping':
            self._serve_json({'status': 'ok', 'timestamp': datetime.now().isoformat()})

        elif path.startswith('/api/stream'):
            from urllib.parse import parse_qs
            params   = parse_qs(urlparse(self.path).query)
            filepath = params.get('file', [''])[0]
            if filepath and os.path.exists(filepath):
                ext  = os.path.splitext(filepath)[1].lower()
                mime = {'.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.ogg': 'audio/ogg',
                        '.mp4': 'video/mp4',  '.webm': 'video/webm', '.m4a': 'audio/mp4',
                        '.flac': 'audio/flac'}.get(ext, 'application/octet-stream')
                size = os.path.getsize(filepath)
                rng  = self.headers.get('Range')
                try:
                    if rng:
                        s, e = rng.replace('bytes=', '').split('-')
                        s = int(s); e = int(e) if e else size - 1
                        self.send_response(206)
                        self.send_header('Content-Type', mime)
                        self.send_header('Content-Range', f'bytes {s}-{e}/{size}')
                        self.send_header('Content-Length', e - s + 1)
                        self.send_header('Accept-Ranges', 'bytes')
                        self.end_headers()
                        with open(filepath, 'rb') as f:
                            f.seek(s); self.wfile.write(f.read(e - s + 1))
                    else:
                        self.send_response(200)
                        self.send_header('Content-Type', mime)
                        self.send_header('Content-Length', size)
                        self.send_header('Accept-Ranges', 'bytes')
                        self.end_headers()
                        with open(filepath, 'rb') as f:
                            self.wfile.write(f.read())
                except Exception:
                    pass
            else:
                self.send_error(404)

        elif path.startswith('/images/'):
            filename = path[8:]
            filepath = os.path.join(os.getcwd(), 'web_images', filename)
            if os.path.exists(filepath):
                self.send_response(200)
                ct = ('image/png' if filename.endswith('.png') else
                      'image/jpeg' if filename.lower().endswith(('.jpg', '.jpeg')) else
                      'image/svg+xml' if filename.endswith('.svg') else
                      'application/octet-stream')
                self.send_header('Content-type', ct)
                self.end_headers()
                with open(filepath, 'rb') as f:
                    self.wfile.write(f.read())
            else:
                self.send_error(404)
        else:
            self.send_error(404)

    # ── OPTIONS (CORS) ────────────────────────────────────────────────────────
    def do_OPTIONS(self):
        self.send_response(204)
        self._add_cors_headers()
        self.end_headers()

    # ── POST ─────────────────────────────────────────────────────────────────
    def do_POST(self):
        path = urlparse(self.path).path
        if path == '/api/send':
            self._handle_send()
        elif path == '/api/voice':
            self._handle_voice()
        elif path == '/api/speak':
            self._handle_speak()
        elif path == '/api/imagine':
            self._handle_imagine()
        elif path == '/api/clear':
            self.nova.root.after(0, lambda: self.nova._new_chat())
            self._serve_json({'status': 'cleared'})
        elif path == '/api/upload':
            self._handle_upload()
        elif path == '/api/tts':
            self._handle_tts_toggle()
        else:
            self.send_error(404)

    # ── Low-level helpers ─────────────────────────────────────────────────────
    def _add_cors_headers(self):
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type, Authorization')
        self.send_header('Access-Control-Max-Age', '86400')

    def _serve_html(self):
        self.send_response(200)
        self.send_header('Content-type', 'text/html; charset=utf-8')
        self.send_header('Cache-Control', 'no-cache, no-store, must-revalidate')
        self.send_header('Content-Security-Policy',
            "default-src 'self'; "
            "script-src 'self' 'unsafe-inline' https://cdn.jsdelivr.net https://cdnjs.cloudflare.com; "
            "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
            "font-src 'self' https://fonts.gstatic.com https://cdn.jsdelivr.net; "
            "img-src 'self' data: blob:; media-src 'self' blob: data:; "
            "connect-src 'self'; frame-ancestors 'none';")
        self.end_headers()
        self.wfile.write(HTML_TEMPLATE.encode('utf-8'))

    def _serve_json(self, data):
        self.send_response(200)
        self.send_header('Content-type', 'application/json')
        self._add_cors_headers()
        self.end_headers()
        def _s(o):
            if hasattr(o, 'strftime'): return o.strftime('%H:%M')
            return str(o)
        self.wfile.write(json.dumps(data, default=_s).encode('utf-8'))

    # ── Whisper ASR ───────────────────────────────────────────────────────────
    def _handle_voice(self):
        try:
            length     = int(self.headers['Content-Length'])
            audio_data = self.rfile.read(length)
            self.nova.log(f"[VOICE] Received {len(audio_data)} bytes")
            if not audio_data:
                self._serve_json({'status': 'ok', 'transcript': ''}); return
            whisper = getattr(self.nova, 'whisper', None)
            if not whisper or not whisper.model_loaded:
                self._serve_json({'error': 'Whisper not ready'}); return
            try:
                from pydub import AudioSegment
                import soundfile as sf, io, numpy as np
                seg = AudioSegment.from_file(io.BytesIO(audio_data))
                seg = seg.set_frame_rate(16000).set_channels(1)
                self.nova.log(f"[VOICE] {seg.duration_seconds:.1f}s")
                wav_io = io.BytesIO(); seg.export(wav_io, format='wav'); wav_io.seek(0)
                samples, _ = sf.read(wav_io, dtype='float32')
                text = whisper.asr_model.transcribe(samples, 16000)
                self.nova.log(f"[VOICE] '{text}'")
                self._serve_json({'status': 'ok', 'transcript': text.strip() if text else ''})
            except ImportError:
                self._serve_json({'error': 'pydub not installed'})
            except Exception as e:
                self._serve_json({'error': str(e)})
        except Exception as e:
            self._serve_json({'error': str(e)})


    # ── Edge TTS → browser MP3 ────────────────────────────────────────────────
    def _handle_speak(self):
        try:
            length = int(self.headers['Content-Length'])
            data = json.loads(self.rfile.read(length).decode('utf-8'))
            text = data.get('text', '').strip()
            if not text or not getattr(self.nova, '_web_tts_on', True) \
                    or getattr(self.nova, '_tts_on', False):
                self.send_response(204);
                self.end_headers();
                return

            import asyncio, edge_tts, re

            math_speech = getattr(self.nova, 'math_speech', None)
            if math_speech:
                # Pass raw text — make_speakable_text extracts math BEFORE
                # stripping markdown, so $...$ delimiters survive intact
                clean = math_speech.make_speakable_text(text, speak_math=True)
            else:
                # No math handler — strip display math first, then inline,
                # using unambiguous non-overlapping patterns
                clean = re.sub(r'\$\$[\s\S]+?\$\$', ' ', text)  # display $$...$$
                clean = re.sub(r'\$[^\$\n]{1,300}?\$', ' ', clean)  # inline $...$
                clean = re.sub(r'\\\[[\s\S]+?\\\]', ' ', clean)  # \[...\]
                clean = re.sub(r'\\\(.+?\\\)', ' ', clean)  # \(...\)
                clean = re.sub(r'```[\s\S]*?```', ' ', clean)  # code blocks
                clean = re.sub(r'`[^`]+`', ' ', clean)  # inline code
                clean = re.sub(r'^#{1,6}\s+', '', clean, flags=re.MULTILINE)
                clean = re.sub(r'\*{1,3}([^*]+)\*{1,3}', r'\1', clean)  # bold/italic
                clean = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', clean)  # links
                clean = re.sub(r'^[-*]\s+', '', clean, flags=re.MULTILINE)
                clean = re.sub(r'[^\w\s\.,!?;:\-\(\)\/\+\=%"\']+', ' ', clean)

            # Shared final cleanup
            clean = re.sub(r'\s+', ' ', clean).strip()
            if not clean:
                self.send_response(204);
                self.end_headers();
                return

            engine_obj = getattr(self.nova, 'tts_engine_combo', None)
            engine_str = engine_obj.get() if engine_obj else 'sapi5'
            if engine_str == 'edge':
                vc = getattr(self.nova, 'edge_voice_combo', None)
                voice = vc.get() if vc else 'en-AU-NatashaNeural'
            else:
                voice = 'en-GB-SoniaNeural'

            async def _run():
                comm = edge_tts.Communicate(clean, voice)
                chunks = []
                async for chunk in comm.stream():
                    if chunk['type'] == 'audio':
                        chunks.append(chunk['data'])
                return b''.join(chunks)

            loop = asyncio.new_event_loop()
            audio = loop.run_until_complete(_run())
            loop.close()

            self.send_response(200)
            self.send_header('Content-Type', 'audio/mpeg')
            self.send_header('Content-Length', len(audio))
            self.end_headers()
            self.wfile.write(audio)
        except Exception as e:
            self.nova.log(f"[SPEAK] {e}")
            self.send_response(204);
            self.end_headers()
    # ── Send / Imagine ────────────────────────────────────────────────────────
    def _handle_send(self):
        length = int(self.headers['Content-Length'])
        data   = json.loads(self.rfile.read(length).decode('utf-8'))
        msg    = data.get('message', '')
        if msg:
            self.nova.root.after(0, lambda: self._send_to_nova(msg))
            self._serve_json({'status': 'processing'})

    def _send_to_nova(self, msg):
        if "I've attached" in msg and "```" in msg:
            try:
                import re
                blocks = re.findall(
                    r'\*\*File: (.+?)\*\*.*?Content:\s*```\r?\n(.*?)\n```', msg, re.DOTALL)
                if blocks:
                    comment = msg.split('User comment:')[-1].strip() if 'User comment:' in msg else ''
                    parts   = [f'Content of "{fn}":\n```\n{fc.strip()}\n```' for fn, fc in blocks]
                    clean   = '\n\n'.join(parts) + '\n\n' + (comment or 'Please summarise this content.')
                    self.nova._append_conv("user", clean)
                    self.nova.conversation_history.append({"role": "user", "content": clean})
                    self.nova._thinking = True
                    threading.Thread(target=self.nova._process_input, args=(clean,), daemon=True).start()
                    return
            except Exception as e:
                self.nova.log(f"[UPLOAD] extraction failed: {e}")
        self.nova._append_conv("user", msg)
        self.nova.conversation_history.append({"role": "user", "content": msg})
        self.nova._thinking = True
        threading.Thread(target=self.nova._process_input, args=(msg,), daemon=True).start()

    def _handle_imagine(self):
        length = int(self.headers['Content-Length'])
        data   = json.loads(self.rfile.read(length).decode('utf-8'))
        msg    = data.get('message', '')
        if msg:
            prompt = (f"You are in CREATIVE IMAGINE MODE.\n"
                      f"Approach this with maximum creativity and lateral thinking.\n"
                      f"User prompt: {msg}")
            self.nova.root.after(0, lambda: self._imagine_to_nova(prompt, msg))
            self._serve_json({'status': 'processing'})

    def _imagine_to_nova(self, prompt, original):
        self.nova._append_conv("user", f"✨ [Holodeck] {original}")
        self.nova.conversation_history.append({"role": "user", "content": f"✨ [Holodeck] {original}"})
        self.nova._thinking = True
        threading.Thread(target=self.nova._process_input, args=(prompt,), daemon=True).start()

    # ── TTS toggle ────────────────────────────────────────────────────────────
    def _handle_tts_toggle(self):
        current = getattr(self.nova, '_web_tts_on', True)
        self.nova._web_tts_on = not current
        if not self.nova._web_tts_on:
            self.nova.root.after(0, self.nova._stop_speaking)
        else:
            self.nova._tts_stop = False
        self.nova.log(f"[WEB TTS] {'Enabled' if self.nova._web_tts_on else 'Disabled'}")
        self._serve_json({'tts': self.nova._web_tts_on})

    # ── File upload (full featured) ───────────────────────────────────────────
    def _handle_upload(self):
        try:
            MAX_FILE    = 100 * 1024 * 1024
            MAX_TOTAL   = 500 * 1024 * 1024
            MAX_PREVIEW = 50_000
            MAX_CODE    = 100_000

            ct  = self.headers.get('Content-Type', '')
            cl  = int(self.headers.get('Content-Length', 0))
            if cl > MAX_TOTAL:
                self.send_error(413, 'Upload too large'); return
            if 'multipart/form-data' not in ct:
                self.send_error(400, 'Expected multipart/form-data'); return

            boundary = None
            for part in ct.split(';'):
                p = part.strip()
                if p.startswith('boundary='):
                    boundary = p[9:].strip().encode('utf-8'); break
            if not boundary:
                self.send_error(400, 'Missing boundary'); return

            raw  = self.rfile.read(cl)
            files = {}
            for part in raw.split(b'--' + boundary)[1:]:
                if part in (b'--\r\n', b'--', b''): continue
                sep = b'\r\n\r\n' if b'\r\n\r\n' in part else b'\n\n'
                if sep not in part: continue
                hdr, body = part.split(sep, 1)
                if body.endswith(b'\r\n'): body = body[:-2]
                disp = {}
                for line in hdr.decode('utf-8', 'replace').splitlines():
                    if 'Content-Disposition' in line:
                        for tok in line.split(';'):
                            tok = tok.strip()
                            if '=' in tok:
                                k, v = tok.split('=', 1)
                                disp[k.strip()] = v.strip().strip('"')
                if disp.get('filename'):
                    files[disp.get('name', '')] = {'filename': disp['filename'], 'data': body}

            if 'files' not in files:
                self.send_error(400, 'No file field'); return

            obj       = files['files']
            items     = obj if isinstance(obj, list) else [obj]
            for fi in items:
                if len(fi['data']) > MAX_FILE:
                    self.send_error(413, f'{fi["filename"]} too large'); return

            upload_dir = os.path.join(os.getcwd(), 'uploads')
            os.makedirs(upload_dir, exist_ok=True)
            results = []

            IMAGE = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}
            CODE  = {'.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.rs', '.go', '.rb', '.php'}
            TEXT  = {'.txt', '.html', '.css', '.json', '.md', '.csv', '.xml', '.yaml', '.yml',
                     '.ini', '.toml', '.bat', '.sh', '.log', '.conf', '.cfg'}

            for fi in items:
                filename  = os.path.basename(fi['filename'])
                file_data = fi['data']
                name, ext = os.path.splitext(filename)
                extl      = ext.lower()

                save_path = os.path.join(upload_dir, filename)
                n = 1
                while os.path.exists(save_path):
                    save_path = os.path.join(upload_dir, f"{name}_{n}{ext}"); n += 1
                with open(save_path, 'wb') as f:
                    f.write(file_data)

                if extl in IMAGE:
                    preview = f"[IMAGE_FILE: {save_path}]"
                elif extl in CODE:
                    try:
                        t = file_data.decode('utf-8', 'replace')
                        preview = t[:MAX_CODE] + ('\n...[truncated]' if len(t) > MAX_CODE else '')
                    except: preview = f"[Cannot decode: {filename}]"
                elif extl in TEXT:
                    try:
                        t = file_data.decode('utf-8', 'replace')
                        preview = t[:MAX_PREVIEW] + ('\n...[truncated]' if len(t) > MAX_PREVIEW else '')
                    except: preview = f"[Cannot decode: {filename}]"
                elif extl in ('.xlsx', '.xls'):
                    try:
                        import pandas as pd, io as _io
                        xl     = pd.ExcelFile(_io.BytesIO(file_data))
                        sheets = xl.sheet_names
                        df     = pd.read_excel(_io.BytesIO(file_data), sheet_name=sheets[0])
                        preview = f"Excel: {len(sheets)} sheet(s): {', '.join(sheets[:5])}\n"
                        preview += f"Shape: {df.shape[0]}×{df.shape[1]}\n{df.head(5).to_string()}"
                        if len(preview) > MAX_PREVIEW: preview = preview[:MAX_PREVIEW] + '\n...[truncated]'
                    except ImportError: preview = "[Excel: pandas not installed]"
                    except Exception as e: preview = f"[Excel error: {e}]"
                elif extl == '.pdf':
                    try:
                        from pypdf import PdfReader
                        import io as _io
                        reader  = PdfReader(_io.BytesIO(file_data))
                        preview = f"PDF: {len(reader.pages)} pages\n\n"
                        total   = 0
                        for i, pg in enumerate(reader.pages[:3]):
                            t = ' '.join((pg.extract_text() or '').split())
                            if total + len(t) > MAX_PREVIEW: break
                            preview += f"--- Page {i+1} ---\n{t}\n\n"; total += len(t)
                    except ImportError: preview = "[PDF: pypdf not installed]"
                    except Exception as e: preview = f"[PDF error: {e}]"
                elif extl in ('.pptx', '.ppt'):
                    try:
                        from pptx import Presentation
                        import io as _io
                        prs     = Presentation(_io.BytesIO(file_data))
                        preview = f"PowerPoint: {len(prs.slides)} slides\n\n"
                        for i, slide in enumerate(prs.slides[:5]):
                            txts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, 'text') and sh.text.strip()]
                            preview += f"--- Slide {i+1} ---\n" + '\n'.join(txts) + '\n\n'
                            if len(preview) > MAX_PREVIEW: preview = preview[:MAX_PREVIEW] + '\n...[truncated]'; break
                    except ImportError: preview = "[PPTX: python-pptx not installed]"
                    except Exception as e: preview = f"[PPTX error: {e}]"
                elif extl in ('.docx', '.doc'):
                    try:
                        from docx import Document
                        import io as _io
                        doc    = Document(_io.BytesIO(file_data))
                        paras  = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                        preview = '\n\n'.join(paras[:50])
                        if len(preview) > MAX_PREVIEW: preview = preview[:MAX_PREVIEW] + '\n...[truncated]'
                    except ImportError: preview = "[DOCX: python-docx not installed]"
                    except Exception as e: preview = f"[DOCX error: {e}]"
                else:
                    preview = f"[Binary: {filename}, {len(file_data)/(1024*1024):.2f} MB]"

                results.append({
                    'filename': os.path.basename(save_path),
                    'original_name': filename,
                    'size': len(file_data),
                    'path': save_path,
                    'preview': preview,
                })
                self.nova.log(f"[UPLOAD] {filename} ({len(file_data)} bytes)")

            self._serve_json({'status': 'success', 'files': results})
        except Exception as e:
            self.nova.log(f"[UPLOAD ERROR] {e}")
            self.send_error(500, str(e))

    def log_message(self, fmt, *args):
        msg = fmt % args
        if any(p in msg for p in ['/api/history', '/api/state', '/api/ping']): return
        if '"GET /api/' in msg and ('200' in msg or '304' in msg): return
        if ('/images/' in msg or '404' in msg or '500' in msg) and self.nova:
            self.nova.log(f"[WEB] {msg}")


# =============================================================================
#  SERVER WRAPPER
# =============================================================================

class NovaWebServer:
    def __init__(self, nova_app, port=8080, bind_all=False):
        self.nova     = nova_app
        self.port     = port
        self.bind_all = bind_all
        self.server   = None
        self.thread   = None

    def start(self):
        if self.server: return
        NovaWebHandler.nova = self.nova
        script_dir = os.path.dirname(os.path.abspath(__file__))
        cert = os.path.join(script_dir, 'cert.pem')
        key  = os.path.join(script_dir, 'key.pem')
        use_ssl = os.path.exists(cert) and os.path.exists(key)
        host    = "0.0.0.0" if self.bind_all else "127.0.0.1"

        if use_ssl:
            import ssl
            ctx = ssl.SSLContext(ssl.PROTOCOL_TLS_SERVER)
            ctx.load_cert_chain(cert, key)
            class SSLTCPServer(socketserver.TCPServer):
                def get_request(self_):
                    ns, fa = self_.socket.accept()
                    return ctx.wrap_socket(ns, server_side=True), fa
            self.server = SSLTCPServer((host, self.port), NovaWebHandler)
            self.nova.log("[WEB] SSL enabled (HTTPS)"); protocol = "https"
        else:
            self.server = socketserver.TCPServer((host, self.port), NovaWebHandler)
            self.nova.log("[WEB] HTTP mode"); protocol = "http"

        self.server.allow_reuse_address = True
        self.thread = threading.Thread(target=self.server.serve_forever, daemon=True)
        self.thread.start()
        ip = self._get_local_ip() if self.bind_all else "127.0.0.1"
        self.nova.log(f"[WEB] {protocol}://{ip}:{self.port}")

    def _get_local_ip(self):
        import socket
        try:
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80)); ip = s.getsockname()[0]; s.close(); return ip
        except: return "localhost"

    def stop(self):
        if self.server:
            self.server.shutdown(); self.server.server_close()
            self.server = None; self.nova.log("[WEB] Server stopped")

    def open_browser(self):
        proto = "https" if self.bind_all else "http"
        ip    = self._get_local_ip() if self.bind_all else "127.0.0.1"
        webbrowser.open(f"{proto}://{ip}:{self.port}")


# =============================================================================
#  HTML TEMPLATE
# =============================================================================

HTML_TEMPLATE = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, viewport-fit=cover, interactive-widget=resizes-content">
<meta name="apple-mobile-web-app-capable" content="yes">
<meta name="apple-mobile-web-app-status-bar-style" content="black-translucent">
<title> NOVA — KNOWLEDGE COMPUTER</title>
<link href="https://fonts.googleapis.com/css2?family=Share+Tech+Mono&family=Orbitron:wght@400;600;700;800;900&display=swap" rel="stylesheet">
<script src="https://cdnjs.cloudflare.com/ajax/libs/marked/9.1.6/marked.min.js"></script>
<script>
window.MathJax = {
  tex: {
    inlineMath: [['$','$'],['\\(','\\)']],
    displayMath: [['$$','$$'],['\\[','\\]']],
    processEscapes: true
  },
  options: { ignoreHtmlClass:'no-mathjax', processHtmlClass:'mathjax-content' }
};
</script>
<script id="MathJax-script" async src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-chtml.js"></script>

<style>
/* ===== LCARS PALETTE ===== */
:root {
  --bg:       #04080F;
  --gold:     #FFB300;
  --coral:    #FF6B35;
  --blue:     #7EB8FF;
  --teal:     #00D4AA;
  --pale:     #D0E8FF;
  --dimteal:  #008C6E;
  --dimblu:   #3C64A0;
  --dimgold:  #B47800;
  --purple:   #A050DC;
  --red:      #C82828;
}

* { box-sizing:border-box; margin:0; padding:0; }
html, body { background:var(--bg); font-family:'Share Tech Mono',monospace;
  overflow:hidden; height:100%; position:fixed; inset:0; }

/* STARFIELD */
#starfield { position:fixed; inset:0; pointer-events:none; z-index:0; }
.star { position:absolute; background:#fff; border-radius:50%; }
@keyframes twinkle {
  0%,100% { opacity:.12; }
  50%      { opacity:.7; }
}

/* APP SHELL */
#app { position:fixed; inset:0; z-index:2; display:flex; flex-direction:column; }

/* ─── HEADER ─── */
.hdr {
  background:linear-gradient(180deg,#08102A 0%,#04080F 100%);
  border-bottom:3px solid var(--gold); padding:10px 18px;
  flex-shrink:0; position:relative;
}
.hdr::before {
  content:''; position:absolute; bottom:-3px; left:0;
  width:220px; height:3px;
  background:linear-gradient(90deg,var(--gold),var(--coral));
}
.name {
  font-family:'Orbitron',monospace; font-size:26px; font-weight:800;
  color:var(--gold); letter-spacing:4px;
  text-shadow:0 0 12px rgba(255,179,0,.55);
}
.date {
  font-size:10px; color:var(--gold); margin-top:3px;
  opacity:0.9;
}
.stardate { font-size:10px; color:var(--teal); margin-top:3px; }
.status-bar { display:flex; gap:8px; margin-top:8px; flex-wrap:wrap; align-items:center; }

/* status pill */
.pill {
  display:flex; align-items:center; gap:6px;
  background:rgba(0,212,170,.07);
  border-left:3px solid var(--coral);
  padding:4px 11px; font-size:10px; font-weight:bold;
  color:var(--pale); white-space:nowrap;
}
.pill.ok  { border-left-color:var(--teal); }
.pill.btn { cursor:pointer; transition:background .15s; }
.pill.btn:hover  { background:rgba(0,212,170,.16); }
.pill.btn:active { transform:scale(.97); }

/* live status dot */
.dot {
  width:8px; height:8px; border-radius:50%; flex-shrink:0;
  background:var(--red); box-shadow:0 0 6px var(--red);
}
.dot.live {
  background:var(--teal); box-shadow:0 0 7px var(--teal);
  animation:dotPulse 2s ease-in-out infinite;
}
@keyframes dotPulse {
  0%,100% { opacity:1;  box-shadow:0 0 6px var(--teal); }
  50%      { opacity:.4; box-shadow:0 0 14px var(--teal); }
}

/* ─── CONVERSATION ─── */
.conv {
  flex:1; overflow-y:auto; padding:16px;
  scroll-behavior:smooth; -webkit-overflow-scrolling:touch;
}
.conv::-webkit-scrollbar { width:5px; }
.conv::-webkit-scrollbar-track { background:#0A1430; }
.conv::-webkit-scrollbar-thumb { background:var(--blue); border-radius:3px; }

/* messages */
.msg { margin-bottom:18px; animation:fadeUp .3s ease-out forwards; opacity:0; }
@keyframes fadeUp {
  from { opacity:0; transform:translateY(14px); }
  to   { opacity:1; transform:translateY(0); }
}
.msg-hdr {
  font-size:9px; font-weight:bold; letter-spacing:1px;
  margin-bottom:4px; display:flex; align-items:center; gap:8px;
}
.msg.user      .msg-hdr { justify-content:flex-end; color:var(--gold); }
.msg.assistant .msg-hdr { color:var(--teal); }
.msg.system    .msg-hdr { justify-content:center; color:var(--coral); }

.bubble {
  max-width:80%; padding:11px 16px; border-radius:6px;
  font-size:13px; line-height:1.65; color:var(--pale); position:relative;
}
.msg.user      .bubble { background:rgba(255,107,53,.10); border-right:3px solid var(--coral); margin-left:auto; }
.msg.assistant .bubble { background:rgba(0,212,170,.05);  border-left:3px solid var(--teal); }
.msg.system    .bubble {
  background:rgba(255,107,53,.07); border:1px solid var(--coral);
  margin:0 auto; font-size:11px; text-align:center; max-width:90%;
}

/* save button */
.save-btn {
  position:absolute; top:6px; right:6px;
  width:27px; height:27px; border-radius:4px;
  background:rgba(0,212,170,.1); border:1px solid rgba(0,212,170,.3);
  color:var(--teal); font-size:12px; cursor:pointer;
  display:flex; align-items:center; justify-content:center;
  opacity:0; transition:opacity .2s;
}
.msg.assistant:hover .save-btn,
.msg.assistant:active .save-btn { opacity:1; }

/* copy-text button on user bubbles */
.copy-txt-btn {
  position:absolute; top:6px; left:6px;
  width:27px; height:27px; border-radius:4px;
  background:rgba(255,107,53,.1); border:1px solid rgba(255,107,53,.3);
  color:var(--coral); font-size:12px; cursor:pointer;
  display:flex; align-items:center; justify-content:center;
  opacity:0; transition:opacity .2s;
}
.msg.user:hover .copy-txt-btn,
.msg.user:active .copy-txt-btn { opacity:1; }

.ref-num { font-size:8px; color:var(--dimgold); margin-top:4px; text-align:right; }

/* ─── MARKDOWN CONTENT in assistant bubble ─── */
.bubble p { margin-bottom:.65em; }
.bubble p:last-child { margin-bottom:0; }
.bubble ul, .bubble ol { margin:.4em 0 .65em 1.4em; }
.bubble li { margin-bottom:.25em; }
.bubble h1,.bubble h2,.bubble h3,.bubble h4 {
  margin:1em 0 .35em; color:var(--gold);
  font-family:'Orbitron',monospace; font-weight:600;
}
.bubble h1:first-child,.bubble h2:first-child,.bubble h3:first-child { margin-top:0; }
.bubble strong { color:#fff; font-weight:600; }
.bubble em { color:var(--blue); }
.bubble blockquote { border-left:2px solid var(--teal); padding-left:10px; margin:.5em 0; color:var(--dimteal); }
.bubble a { color:var(--blue); text-decoration:none; }
.bubble a:hover { color:var(--teal); text-decoration:underline; }
.bubble hr { border:none; height:1px; margin:1em 0; background:linear-gradient(90deg,transparent,var(--gold),transparent); }
.bubble table { border-collapse:collapse; width:100%; margin:.65em 0; font-size:11px; }
.bubble th,.bubble td { padding:5px 9px; border:1px solid rgba(126,184,255,.2); text-align:left; }
.bubble th { background:rgba(255,179,0,.1); color:var(--gold); font-weight:600; }
.bubble tr:nth-child(even) td { background:rgba(255,255,255,.02); }
.bubble code:not(pre code) {
  background:rgba(0,212,170,.12); color:var(--teal);
  padding:1px 5px; border-radius:3px; font-size:.88em; font-family:monospace;
}

/* ─── CODE BLOCKS ─── */
.code-wrap {
  margin:10px 0; border-radius:7px; overflow:hidden;
  border:1px solid rgba(0,212,170,.18); background:#060C18;
}
.code-hdr {
  display:flex; align-items:center; justify-content:space-between;
  padding:6px 12px; background:rgba(0,0,0,.35);
  border-bottom:1px solid rgba(0,212,170,.12);
}
.tl-wrap { display:flex; gap:5px; }
.tl { width:10px; height:10px; border-radius:50%; }
.tl-r{background:#FF5F57;} .tl-y{background:#FEBC2E;} .tl-g{background:#28C840;}
.cp-btn {
  opacity:0; padding:3px 10px;
  background:rgba(255,255,255,.05);
  border:1px solid rgba(126,184,255,.3); border-radius:999px;
  color:var(--blue); font-size:10px; cursor:pointer; transition:opacity .2s;
  font-family:'Share Tech Mono',monospace;
}
.code-wrap:hover .cp-btn,
.code-wrap:active .cp-btn { opacity:1; }
pre { padding:14px; overflow-x:auto; font-size:12px; margin:0; }

/* ─── TYPING INDICATOR ─── */
#typing-ind {
  display:none; padding:10px 18px;
  background:rgba(0,212,170,.04); border-left:3px solid var(--teal);
  align-items:center; gap:14px; font-size:10px; color:var(--teal); flex-shrink:0;
}
#typing-ind.active { display:flex; }
.warp-bars { display:flex; gap:5px; align-items:center; }
.warp-bar { width:5px; background:var(--gold); animation:warpPulse .8s ease-in-out infinite; }
@keyframes warpPulse {
  0%,100% { height:6px;  opacity:.35; }
  50%      { height:20px; opacity:1; }
}

/* ─── CONSOLE ─── */
.console {
  flex-shrink:0;
  background:linear-gradient(0deg,#060C18 0%,#0A1430 100%);
  border-top:3px solid var(--gold); padding:12px 18px; position:relative;
}
.console::before {
  content:''; position:absolute; top:-3px; right:0;
  width:200px; height:3px;
  background:linear-gradient(270deg,var(--gold),var(--coral));
}
/* drop zone */
.drop-zone {
  border:1px dashed rgba(126,184,255,.22); border-radius:5px;
  background:rgba(4,8,15,.7); padding:7px 14px; margin-bottom:9px;
  cursor:pointer; display:flex; align-items:center; gap:8px;
  font-size:10px; color:var(--dimteal); transition:all .25s;
}
.drop-zone.drag-over { border-color:var(--teal); background:rgba(0,212,170,.08); color:var(--teal); }
.file-prev { display:flex; flex-wrap:wrap; gap:6px; margin-bottom:8px; }
.fpi {
  background:rgba(126,184,255,.1); border-radius:4px;
  padding:3px 9px; font-size:10px; color:var(--pale);
  display:flex; align-items:center; gap:5px;
}
.fpi button { background:none; border:none; color:var(--coral); cursor:pointer; font-size:11px; padding:0 2px; }

.input-row { display:flex; gap:9px; align-items:flex-end; }
.ifield {
  flex:1; background:rgba(4,8,15,.9); border:1px solid var(--dimblu);
  border-radius:5px; padding:13px 15px; color:var(--pale);
  font-family:'Share Tech Mono',monospace; font-size:13px;
  resize:none; outline:none; transition:border-color .2s;
  min-height:48px; max-height:130px;
}
.ifield:focus { border-color:var(--blue); box-shadow:0 0 8px rgba(126,184,255,.2); }
.ifield::placeholder { color:var(--dimteal); }

.btn {
  padding:11px 18px; border:none; border-radius:5px;
  font-family:'Orbitron',monospace; font-weight:700;
  font-size:11px; cursor:pointer; letter-spacing:1px;
  transition:transform .15s; white-space:nowrap;
  -webkit-tap-highlight-color:transparent;
}
.btn:active { transform:scale(.96); }
.btn-tx  { background:linear-gradient(135deg,var(--coral),#CC4422); color:#fff; box-shadow:0 0 10px rgba(255,107,53,.35); }
.btn-hd  { background:linear-gradient(135deg,var(--purple),#6B2FA0); color:#fff; box-shadow:0 0 10px rgba(160,80,220,.35); }
.btn-mic { background:linear-gradient(135deg,var(--teal),#008C6E);   color:#fff; box-shadow:0 0 10px rgba(0,212,170,.35); }
.btn-mic.recording { background:linear-gradient(135deg,var(--red),#8C1010); animation:micPulse 1s ease-in-out infinite; }
@keyframes micPulse {
  0%,100% { box-shadow:0 0 10px rgba(200,40,40,.35); }
  50%      { box-shadow:0 0 22px rgba(200,40,40,.75); }
}
.console-foot { display:flex; justify-content:space-between; margin-top:6px; font-size:9px; color:var(--dimteal); }

/* ─── MEDIA ─── */
.media-wrap { margin:10px 0; text-align:center; }
.media-wrap img { max-width:100%; height:auto; border-radius:7px; cursor:pointer;
  box-shadow:0 2px 8px rgba(0,0,0,.4); display:block; margin:0 auto; }
.media-cap { font-size:9px; color:var(--dimteal); margin-top:4px; }

/* ─── TOAST ─── */
.toast {
  position:fixed; bottom:18px; left:50%; transform:translateX(-50%);
  padding:9px 18px; border-radius:4px; font-size:11px; font-weight:bold;
  z-index:10000; animation:slideUp .3s ease-out; white-space:nowrap;
  background:rgba(0,212,170,.95); color:var(--bg);
  font-family:'Share Tech Mono',monospace;
}
@keyframes slideUp {
  from { transform:translateX(-50%) translateY(100%); opacity:0; }
  to   { transform:translateX(-50%) translateY(0);    opacity:1; }
}

/* ─── WELCOME ─── */
.welcome { text-align:center; padding:50px 20px; animation:fadeUp .5s ease-out forwards; opacity:0; }
.delta   { font-size:70px; margin-bottom:16px; animation:dPulse 3s ease-in-out infinite; }
@keyframes dPulse { 0%,100%{opacity:.8;text-shadow:0 0 20px var(--gold);}50%{opacity:1;text-shadow:0 0 40px var(--gold);} }
.w-title { font-family:'Orbitron',monospace; font-size:30px; font-weight:800; color:var(--gold); letter-spacing:8px; margin-bottom:8px; }
.w-sub   { font-size:12px; color:var(--teal); margin-bottom:30px; }
.chip-grid { display:grid; grid-template-columns:repeat(auto-fit,minmax(180px,1fr)); gap:10px; max-width:800px; margin:0 auto; }
.chip {
  background:rgba(126,184,255,.06); border-left:3px solid var(--blue);
  padding:9px 14px; font-size:11px; color:var(--pale); cursor:pointer;
  transition:all .2s; text-align:left; font-family:'Share Tech Mono',monospace;
}
.chip:hover { background:rgba(126,184,255,.13); transform:translateX(4px); }

/* ─── RESPONSIVE ─── */
@media (max-width:600px) {
  .ship-name { font-size:18px; letter-spacing:2px; }
  .bubble { max-width:92%; font-size:12px; }
  .btn { padding:10px 12px; font-size:10px; }
  .input-row { flex-wrap:wrap; }
  .ifield { width:100%; font-size:16px; min-height:70px; }
}
</style>
</head>
<body>

<div id="starfield"></div>

<div id="app">

  <!-- ══ HEADER ══ -->
  <div class="hdr">
    <div class="name">NOVA</div>
    <div class="date" id="date">DEEP KNOWLEDGE </div>
    <div class="status-bar">
      <!-- online dot pill -->
      <div class="pill ok" id="conn-pill">
        <span class="dot" id="status-dot"></span>
        <span id="conn-lbl">CONNECTING...</span>
      </div>
      <div class="pill ok"><span>⚡</span> KNOWLEDGE CORE</div>
      <div class="pill ok"><span>📡</span> ENTER QUESTION</div>
      <!-- TTS mute toggle -->
      <div class="pill btn" id="tts-pill" onclick="toggleTTS()">
        <span id="tts-icon">🔊</span><span id="tts-lbl">&nbsp;TTS ON</span>
      </div>
      <!-- clear logs -->
      <div class="pill btn" onclick="clearChat()">🗑&nbsp;CLEAR LOGS</div>
    </div>
  </div>

  <!-- ══ CONVERSATION ══ -->
  <div class="conv" id="conv">
    <div class="welcome" id="welcome">
      <div class="delta">🛡️</div>
      <div class="w-title">COMPUTER ONLINE</div>
      <div class="w-sub">AWAITING  INQUIRY — ALL SYSTEMS GO</div>
      <div class="chip-grid">
        <div class="chip" onclick="useSuggestion(this)">⚡ Explain quantum entanglement</div>
        <div class="chip" onclick="useSuggestion(this)">💻 Write a Python async example</div>
        <div class="chip" onclick="useSuggestion(this)">🕳️ How do black holes form?</div>
        <div class="chip" onclick="useSuggestion(this)">⚙️ Debug my NOVA code</div>
        <div class="chip" onclick="useSuggestion(this)">💡 Creative story ideas</div>
        <div class="chip" onclick="useSuggestion(this)">📊 Analyse this dataset</div>
      </div>
    </div>
  </div>

  <!-- ══ TYPING INDICATOR ══ -->
  <div id="typing-ind">
    <div class="warp-bars">
      <div class="warp-bar" style="animation-delay:0s"></div>
      <div class="warp-bar" style="animation-delay:.1s"></div>
      <div class="warp-bar" style="animation-delay:.2s"></div>
      <div class="warp-bar" style="animation-delay:.3s"></div>
      <div class="warp-bar" style="animation-delay:.4s"></div>
    </div>
    <span>COMPUTER PROCESSING INQUIRY...</span>
  </div>

  <!-- ══ CONSOLE ══ -->
  <div class="console">
    <div class="drop-zone" id="drop-zone">
      <span>📎</span><span>DROP FILES OR CLICK TO UPLOAD</span>
      <input type="file" id="file-input" multiple style="display:none">
    </div>
    <div class="file-prev" id="file-prev"></div>
    <div class="input-row">
      <textarea class="ifield" id="msg-in" rows="1"
        placeholder="ENTER INQUIRY FOR COMPUTER..."></textarea>
      <button class="btn btn-mic" id="mic-btn" onclick="toggleMic()">🎤 VOICE</button>
      <button class="btn btn-hd" onclick="sendImagine()">✨ AESTHETIC</button>
      <button class="btn btn-tx" onclick="sendMessage()">📡 ENTER</button>
    </div>
    <div class="console-foot">
      <span>SHIFT+ENTER: NEW LINE &nbsp;·&nbsp; ENTER: TRANSMIT &nbsp;·&nbsp; DROP FILES TO UPLOAD</span>
      <span id="char-info"></span>
    </div>
  </div>

</div><!-- #app -->

<script>
// ═══════════════════════════════════════════════════════════
//  STARFIELD
// ═══════════════════════════════════════════════════════════
(function(){
  const f=document.getElementById('starfield');
  for(let i=0;i<300;i++){
    const s=document.createElement('div'); s.className='star';
    const sz=Math.random()*2+.4;
    s.style.cssText=`width:${sz}px;height:${sz}px;left:${Math.random()*100}%;`+
      `top:${Math.random()*100}%;opacity:${.12+Math.random()*.5};`+
      `animation:twinkle ${2+Math.random()*3}s ease-in-out ${-(Math.random()*5).toFixed(1)}s infinite;`;
    f.appendChild(s);
  }
})();

// ═══════════════════════════════════════════════════════════
//  STARDATE
// ═══════════════════════════════════════════════════════════
const _T0=Date.now();
function tickStardate(){
  const el=document.getElementById('stardate');
  if(el) el.textContent='DEEP KNOWLEDGE DIVISION — STARDATE '+
    (47634.2+(Date.now()-_T0)/1e6).toFixed(1);
}
tickStardate(); setInterval(tickStardate,1000);

// ═══════════════════════════════════════════════════════════
//  STATE
// ═══════════════════════════════════════════════════════════
const conv       = document.getElementById('conv');
const msgIn      = document.getElementById('msg-in');
let lastCount    = 0, lastContent = '', isThinking = false, userScrolled = false;
let audioCtx     = null, curAudio = null;
let mediaRec     = null, audioChunks = [], isRec = false;
let pendingFiles = [], refCounter = 47291;

// ═══════════════════════════════════════════════════════════
//  TOAST
// ═══════════════════════════════════════════════════════════
function toast(msg, err=false){
  const t=document.createElement('div'); t.className='toast';
  t.style.background=err?'rgba(200,40,40,.95)':'rgba(0,212,170,.95)';
  t.style.color=err?'#fff':'#04080F';
  t.textContent=msg; document.body.appendChild(t);
  setTimeout(()=>t.remove(),3000);
}

// ═══════════════════════════════════════════════════════════
//  ONLINE STATUS DOT  (health-check every 4 s)
// ═══════════════════════════════════════════════════════════
async function checkHealth(){
  const dot=document.getElementById('status-dot');
  const lbl=document.getElementById('conn-lbl');
  try{
    const r=await fetch('/api/ping');
    if(r.ok){ dot.classList.add('live'); lbl.textContent='ONLINE'; }
    else     { dot.classList.remove('live'); lbl.textContent='DEGRADED'; }
  }catch{    dot.classList.remove('live'); lbl.textContent='OFFLINE'; }
}
checkHealth(); setInterval(checkHealth,4000);

// ═══════════════════════════════════════════════════════════
//  WEB TTS  (edge_tts → MP3 → Web Audio)
// ═══════════════════════════════════════════════════════════
function unlockAudio(){
  if(audioCtx) return;
  audioCtx=new(window.AudioContext||window.webkitAudioContext)();
  const buf=audioCtx.createBuffer(1,1,22050);
  const src=audioCtx.createBufferSource();
  src.buffer=buf; src.connect(audioCtx.destination); src.start(0); audioCtx.resume();
}
document.addEventListener('touchstart',unlockAudio,{once:true});
document.addEventListener('click',     unlockAudio,{once:true});

function stopAudio(){
  if(curAudio){ try{curAudio.stop();}catch(e){} curAudio=null; }
}

async function playTTS(text){
  stopAudio();
  try{
    const r=await fetch('/api/speak',{method:'POST',
      headers:{'Content-Type':'application/json'},body:JSON.stringify({text})});
    if(!r.ok||r.status===204) return;
    const buf=await r.arrayBuffer();
    if(!buf.byteLength) return;
    if(!audioCtx||audioCtx.state==='closed')
      audioCtx=new(window.AudioContext||window.webkitAudioContext)();
    await audioCtx.resume();
    const decoded=await audioCtx.decodeAudioData(buf);
    const src=audioCtx.createBufferSource();
    curAudio=src; src.buffer=decoded; src.connect(audioCtx.destination); src.start(0);
    src.onended=()=>{ curAudio=null; };
  }catch(e){ console.warn('TTS:',e); curAudio=null; }
}

async function toggleTTS(){
  const r=await fetch('/api/tts',{method:'POST'});
  const d=await r.json(); const on=d.tts;
  document.getElementById('tts-icon').textContent=on?'🔊':'🔇';
  document.getElementById('tts-lbl').textContent=on?' TTS ON':' TTS OFF';
  document.getElementById('tts-pill').style.borderLeftColor=on?'var(--teal)':'var(--coral)';
  if(!on) stopAudio();
  toast(on?'🔊 TTS enabled':'🔇 TTS muted');
}

async function syncTTS(){
  try{
    const r=await fetch('/api/state'); const d=await r.json();
    const on=(d.web_tts!==undefined)?d.web_tts:true;
    document.getElementById('tts-icon').textContent=on?'🔊':'🔇';
    document.getElementById('tts-lbl').textContent=on?' TTS ON':' TTS OFF';
    document.getElementById('tts-pill').style.borderLeftColor=on?'var(--teal)':'var(--coral)';
  }catch(e){}
}
syncTTS();

// ═══════════════════════════════════════════════════════════
//  VOICE INPUT
// ═══════════════════════════════════════════════════════════
function toggleMic(){
  if(isRec){ stopMic(); return; }
  navigator.mediaDevices.getUserMedia({audio:true})
    .then(stream=>{
      audioChunks=[]; mediaRec=new MediaRecorder(stream);
      mediaRec.ondataavailable=e=>{ if(e.data.size>0) audioChunks.push(e.data); };
      mediaRec.onstop=async()=>{
        stream.getTracks().forEach(t=>t.stop());
        await sendVoice(new Blob(audioChunks,{type:'audio/webm'}));
      };
      mediaRec.start(); isRec=true;
      const b=document.getElementById('mic-btn');
      b.classList.add('recording'); b.textContent='⏹ STOP';
      toast('🎤 Recording… press again to stop');
    })
    .catch(()=>toast('❌ Microphone access denied',true));
}
function stopMic(){
  if(mediaRec&&mediaRec.state!=='inactive') mediaRec.stop();
  isRec=false;
  const b=document.getElementById('mic-btn');
  b.classList.remove('recording'); b.textContent='🎤 VOICE';
}
async function sendVoice(blob){
  toast('⏳ Transcribing...');
  try{
    const r=await fetch('/api/voice',{method:'POST',
      headers:{'Content-Type':'audio/webm','Content-Length':blob.size},body:blob});
    const d=await r.json();
    if(d.transcript){
      msgIn.value=d.transcript; msgIn.dispatchEvent(new Event('input'));
      toast('✓ Transcript ready');
    } else toast('⚠ No speech detected',true);
  }catch{ toast('❌ Voice error',true); }
}

// ═══════════════════════════════════════════════════════════
//  CODE-BLOCK COPY BUTTONS
// ═══════════════════════════════════════════════════════════
function addCodeCopy(container){
  container.querySelectorAll('pre').forEach(pre=>{
    if(pre.closest('.code-wrap')) return;
    const wrap=document.createElement('div'); wrap.className='code-wrap';
    pre.parentNode.insertBefore(wrap,pre); wrap.appendChild(pre);
    const hdr=document.createElement('div'); hdr.className='code-hdr';
    hdr.innerHTML='<div class="tl-wrap"><div class="tl tl-r"></div><div class="tl tl-y"></div><div class="tl tl-g"></div></div>';
    const cb=document.createElement('button'); cb.className='cp-btn'; cb.textContent='📋 COPY';
    cb.onclick=e=>{ e.stopPropagation();
      navigator.clipboard.writeText(pre.textContent).then(()=>{
        cb.textContent='✓ COPIED'; setTimeout(()=>cb.textContent='📋 COPY',2000);
        toast('✓ Code copied');
      });
    };
    hdr.appendChild(cb); wrap.insertBefore(hdr,pre);
  });
}

// ═══════════════════════════════════════════════════════════
//  MATH RENDERING
// ═══════════════════════════════════════════════════════════
async function renderMath(el){
  if(!window.MathJax||!el) return;
  try{ MathJax.typesetClear([el]); await MathJax.typesetPromise([el]); }
  catch(e){ console.warn('MathJax:',e); }
}

// ═══════════════════════════════════════════════════════════
//  VIDEO PLAYER
// ═══════════════════════════════════════════════════════════
function buildVideo(url,name){
  const wrap=document.createElement('div');
  wrap.style.cssText='background:#060C18;border-radius:6px;overflow:hidden;margin-top:8px;border:1px solid rgba(0,212,170,.2);';
  const vid=document.createElement('video');
  vid.style.cssText='width:100%;max-height:340px;display:block;background:#000;';
  vid.appendChild(Object.assign(document.createElement('source'),{src:url}));
  wrap.appendChild(vid);

  const ctrl=document.createElement('div');
  ctrl.style.cssText='padding:7px 10px;background:rgba(0,0,0,.5);';

  // progress bar
  const pw=document.createElement('div');
  pw.style.cssText='height:4px;background:rgba(255,255,255,.15);border-radius:2px;margin-bottom:7px;cursor:pointer;';
  const pb=document.createElement('div');
  pb.style.cssText='height:100%;width:0%;background:linear-gradient(90deg,var(--gold),var(--coral));border-radius:2px;pointer-events:none;';
  pw.appendChild(pb);
  pw.onclick=e=>{ const r=pw.getBoundingClientRect(); vid.currentTime=((e.clientX-r.left)/r.width)*vid.duration; };
  ctrl.appendChild(pw);

  const row=document.createElement('div'); row.style.cssText='display:flex;gap:6px;align-items:center;flex-wrap:wrap;';
  const bs='border-radius:3px;padding:3px 8px;cursor:pointer;font-size:10px;border:1px solid rgba(0,212,170,.25);background:rgba(0,212,170,.08);color:var(--pale);font-family:"Share Tech Mono",monospace;';

  const playB=document.createElement('button'); playB.style.cssText=bs; playB.textContent='▶/⏸';
  playB.onclick=()=>vid.paused?vid.play():vid.pause(); row.appendChild(playB);
  [[-10,'⏪-10s'],[10,'10s⏩']].forEach(([d,l])=>{
    const b=document.createElement('button'); b.style.cssText=bs; b.textContent=l;
    b.onclick=()=>vid.currentTime+=d; row.appendChild(b);
  });
  const tim=document.createElement('span');
  tim.style.cssText='color:var(--dimteal);font-size:10px;min-width:80px;';
  tim.textContent='0:00/0:00'; row.appendChild(tim);

  const spd=document.createElement('span'); spd.style.cssText='color:var(--dimteal);font-size:10px;'; spd.textContent='SPD:'; row.appendChild(spd);
  const spdBtns=[];
  [0.5,1,1.5,2].forEach(s=>{
    const b=document.createElement('button'); b.style.cssText=bs; b.textContent=s+'x';
    b.onclick=()=>{ vid.playbackRate=s; spdBtns.forEach(x=>x.style.background='rgba(0,212,170,.08)'); b.style.background='rgba(0,212,170,.3)'; };
    if(s===1) b.style.background='rgba(0,212,170,.3)';
    spdBtns.push(b); row.appendChild(b);
  });
  const fsB=document.createElement('button'); fsB.style.cssText=bs+'margin-left:auto;'; fsB.textContent='⛶ FULL';
  fsB.onclick=()=>vid.requestFullscreen(); row.appendChild(fsB);

  const fmt=s=>Math.floor(s/60)+':'+String(Math.floor(s%60)).padStart(2,'0');
  vid.addEventListener('timeupdate',()=>{
    if(vid.duration){ pb.style.width=(vid.currentTime/vid.duration*100)+'%'; tim.textContent=fmt(vid.currentTime)+'/'+fmt(vid.duration); }
  });

  ctrl.appendChild(row);
  const cap=document.createElement('div'); cap.className='media-cap'; cap.textContent='🎬 '+name;
  ctrl.appendChild(cap); wrap.appendChild(ctrl);
  return wrap;
}

// ═══════════════════════════════════════════════════════════
//  SAVE RESPONSE
// ═══════════════════════════════════════════════════════════
const SAVECSS=`
*{box-sizing:border-box;}
body{background:#04080F;color:#D0E8FF;font-family:'Share Tech Mono',monospace;padding:40px 20px;line-height:1.75;}
.c{max-width:960px;margin:0 auto;background:rgba(10,20,40,.75);border:1px solid rgba(0,212,170,.25);border-radius:10px;padding:32px;}
h1,h2,h3,h4{font-family:'Orbitron',monospace;color:#FFB300;margin:1em 0 .4em;}
h1{font-size:1.5em;} h2{font-size:1.2em;} h3{font-size:1em;}
h1:first-child,h2:first-child{margin-top:0;}
.ts{color:#3C64A0;font-size:10px;margin-bottom:24px;}
p{margin-bottom:.75em;}
hr{border:none;height:1px;background:linear-gradient(90deg,transparent,#FFB300,transparent);margin:1.2em 0;}
pre{background:#060C18;padding:14px;border-radius:6px;overflow-x:auto;border:1px solid rgba(0,212,170,.2);margin:.75em 0;}
code{background:rgba(0,212,170,.12);color:#00D4AA;padding:1px 5px;border-radius:3px;font-size:.9em;}
blockquote{border-left:3px solid #00D4AA;padding-left:12px;margin:.75em 0;color:#008C6E;}
a{color:#7EB8FF;}
ul,ol{margin:.5em 0 .75em 1.4em;}
li{margin-bottom:.25em;}
strong{color:#fff;}
table{border-collapse:collapse;width:100%;margin:.85em 0;font-size:13px;overflow-x:auto;display:block;}
thead tr{background:rgba(255,179,0,.12);}
th{padding:7px 11px;border:1px solid rgba(126,184,255,.25);color:#FFB300;font-weight:bold;text-align:left;}
td{padding:6px 11px;border:1px solid rgba(126,184,255,.18);text-align:left;}
tbody tr:nth-child(even){background:rgba(255,255,255,.03);}
img{max-width:100%;border-radius:6px;}
.pb{position:fixed;top:14px;right:14px;padding:8px 18px;
background:linear-gradient(135deg,#FF6B35,#CC4422);border:none;border-radius:4px;
color:#fff;cursor:pointer;font-family:'Orbitron',monospace;font-weight:bold;font-size:11px;z-index:999;}
@media print{.pb{display:none;}body{background:#fff;color:#000;}.c{background:#fff;border:none;}}
`;
function saveResponse(html){
  const action=confirm('Save response?\n\nOK = download HTML file\nCancel = open new tab');
  const full=`<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8">
<title>Nova Response — ${new Date().toLocaleString()}</title>
<link href="https://fonts.googleapis.com/css2?family=Orbitron:wght@700&family=Share+Tech+Mono&display=swap" rel="stylesheet">
<script>window.MathJax={tex:{inlineMath:[['$','$'],['\\\\(','\\\\)']],displayMath:[['$$','$$'],['\\\\[','\\\\]']],processEscapes:true},options:{skipHtmlTags:['script','noscript','style','textarea','pre']}};<\/script>
<script src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-chtml.js"><\/script>
<style>${SAVECSS}</style></head><body>
<button class="pb" onclick="window.print()">🖨 PRINT / PDF</button>
<div class="c">
<h1>⚡ NOVA — KNOWLEDGE COMPUTER</h1>
<div class="ts">Generated: ${new Date().toLocaleString()}</div>
${html}
</div></body></html>`;
  if(action){
    const a=document.createElement('a');
    a.href=URL.createObjectURL(new Blob([full],{type:'text/html'}));
    a.download=`nova_${Date.now()}.html`; a.click(); URL.revokeObjectURL(a.href);
    toast('✓ Response saved');
  } else {
    const w=window.open(); w.document.write(full); w.document.close();
    toast('✓ Opened in new tab');
  }
}

// ═══════════════════════════════════════════════════════════
//  ADD MESSAGE
// ═══════════════════════════════════════════════════════════
function addMessage(role, content, scroll=true){
  const welcome=document.getElementById('welcome');
  if(welcome && role!=='system') welcome.remove();

  refCounter++;
  const ref='REF-'+refCounter+'-'+['ALPHA','BETA','GAMMA'][refCounter%3];
  const time=new Date().toLocaleTimeString('en-US',{hour:'2-digit',minute:'2-digit'});

  const row=document.createElement('div'); row.className='msg '+role;
  const label={user:'HUMAN',assistant:'COMPUTER',system:'SYSTEM'}[role]||role.toUpperCase();

  row.innerHTML=
    `<div class="msg-hdr"><span>${label}</span><span style="opacity:.55;font-size:8px">${time}</span></div>`+
    `<div class="bubble"></div>`+
    (role==='assistant'?`<div class="ref-num">${ref}</div>`:'');

  const bubble=row.querySelector('.bubble');

  // Save button (assistant)
  // We re-render from the raw `content` string (with $...$ intact) rather than
  // cloning the live DOM, so the saved file's MathJax can render equations fresh.
  if(role==='assistant'){
    const sb=document.createElement('button'); sb.className='save-btn'; sb.title='Save response'; sb.textContent='💾';
    sb.onclick=()=>{
      let saveHtml;
      if(window.marked){
        const maths=[];
        let s=content;
        const ph=m=>{ maths.push(m); return '@@M'+(maths.length-1)+'@@'; };
        s=s.replace(/\$\$[\s\S]+?\$\$/g,ph);
        s=s.replace(/\$[^\$\n]+?\$/g,ph);
        let html=marked.parse(s);
        html=html.replace(/@@M(\d+)@@/g,(_,i)=>maths[+i]);
        saveHtml=html;
      } else {
        saveHtml=content.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/\n/g,'<br>');
      }
      saveResponse(saveHtml);
    };
    bubble.appendChild(sb);
  }

  // Copy-text button (user)
  if(role==='user'){
    const cb=document.createElement('button'); cb.className='copy-txt-btn'; cb.title='Copy text'; cb.textContent='📋';
    cb.onclick=()=>{ navigator.clipboard.writeText(content).then(()=>toast('✓ Copied')); };
    bubble.appendChild(cb);
  }

  const trim=(content||'').trim();

  // ── special message types ────────────────────────────────
  if(trim.startsWith('[IMAGE:')){
    const fn=trim.slice(7,-1);
    const w=document.createElement('div'); w.className='media-wrap';
    const img=document.createElement('img'); img.src='/images/'+fn;
    img.onclick=()=>window.open('/images/'+fn);
    img.onload=()=>{ if(!userScrolled) conv.scrollTop=conv.scrollHeight; };
    const cap=document.createElement('div'); cap.className='media-cap'; cap.textContent='📷 Click to enlarge';
    w.appendChild(img); w.appendChild(cap); bubble.appendChild(w);

  } else if(trim.startsWith('[DIAGRAM:')){
    const fn=trim.slice(9,-1);
    const w=document.createElement('div'); w.className='media-wrap';
    const img=document.createElement('img'); img.src='/images/'+fn;
    img.onclick=()=>window.open('/images/'+fn);
    img.onload=()=>{ if(!userScrolled) conv.scrollTop=conv.scrollHeight; };
    const cap=document.createElement('div'); cap.className='media-cap'; cap.textContent='📊 Generated Diagram';
    w.appendChild(img); w.appendChild(cap); bubble.appendChild(w);

  } else if(content && content.includes('[AUDIO:')){
    const m=content.match(/\[AUDIO:(.+?)\]/);
    if(m){
      const url='/api/stream?file='+encodeURIComponent(m[1]);
      const name=m[1].split(/[/\\]/).pop();
      const txt=document.createElement('div'); txt.textContent=content.replace(/\[AUDIO:.+?\]/,'').trim();
      bubble.appendChild(txt);
      const w=document.createElement('div'); w.className='media-wrap';
      const aud=document.createElement('audio'); aud.controls=true; aud.style.width='100%';
      aud.appendChild(Object.assign(document.createElement('source'),{src:url}));
      const cap=document.createElement('div'); cap.className='media-cap'; cap.textContent='🎵 '+name;
      w.appendChild(aud); w.appendChild(cap); bubble.appendChild(w);
    }

  } else if(content && content.includes('[VIDEO:')){
    const m=content.match(/\[VIDEO:(.+?)\]/);
    if(m){
      const url='/api/stream?file='+encodeURIComponent(m[1]);
      const name=m[1].split(/[/\\]/).pop();
      const txt=document.createElement('div'); txt.textContent=content.replace(/\[VIDEO:.+?\]/,'').trim();
      bubble.appendChild(txt); bubble.appendChild(buildVideo(url,name));
    }

  } else {
    // Markdown + Math — always append into a child div so the save-btn
    // DOM node (with its closure-based onclick) is never serialised/destroyed.
    const contentDiv = document.createElement('div');
    if(window.marked && role==='assistant'){
      const maths=[];
      let s=content;
      const ph=m=>{ maths.push(m); return '@@M'+(maths.length-1)+'@@'; };
      s=s.replace(/\$\$[\s\S]+?\$\$/g,ph);
      s=s.replace(/\$[^\$\n]+?\$/g,ph);
      let html=marked.parse(s);
      html=html.replace(/@@M(\d+)@@/g,(_,i)=>maths[+i]);
      contentDiv.innerHTML = html;
    } else {
      contentDiv.innerHTML = content
        .replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;')
        .replace(/\n/g,'<br>');
    }
    bubble.appendChild(contentDiv);
    addCodeCopy(bubble);
    if(window.MathJax && role==='assistant' &&
       (contentDiv.innerHTML.includes('$') || contentDiv.innerHTML.includes('\\(')))
      setTimeout(()=>renderMath(bubble),80);
  }

  conv.appendChild(row);
  if(scroll && !userScrolled) row.scrollIntoView({behavior:'smooth',block:'end'});
}

// ═══════════════════════════════════════════════════════════
//  TYPING / STATUS
// ═══════════════════════════════════════════════════════════
function setThinking(on){
  document.getElementById('typing-ind').classList.toggle('active',on);
  const lbl=document.getElementById('conn-lbl');
  if(on) lbl.textContent='PROCESSING...'; else lbl.textContent='ONLINE';
}

// ═══════════════════════════════════════════════════════════
//  POLL FOR RESPONSE
// ═══════════════════════════════════════════════════════════
function startPoll(){
  let wait=0, sawThinking=false, finished=false;
  const iv=setInterval(async()=>{
    if(finished) return;
    wait++;
    try{
      const [hR,sR]=await Promise.all([fetch('/api/history'),fetch('/api/state')]);
      const hist=await hR.json(); const state=await sR.json();
      if(state.thinking) sawThinking=true;

      if(hist.length>lastCount){
        // New message arrived
        finished=true; clearInterval(iv);
        setThinking(false); isThinking=false; loadHistory();

      } else if(sawThinking&&!state.thinking&&wait>4){
        // Nova finished — give history 2s to flush, background loadHistory catches it
        finished=true; clearInterval(iv);
        setTimeout(()=>{ setThinking(false); isThinking=false; loadHistory(); },2000);

      } else if(wait>1200){
        // 10-min silent safety net — no error message, background loadHistory
        // (running every 1500ms) will still deliver the reply when it arrives
        finished=true; clearInterval(iv);
        setThinking(false); isThinking=false;
      }
    }catch(e){}
  },500);
}

// ═══════════════════════════════════════════════════════════
//  SEND MESSAGE / IMAGINE
// ═══════════════════════════════════════════════════════════
async function sendMessage(){
  if(isThinking) return;
  const raw=msgIn.value.trim();
  const payload=buildFilePayload();
  if(!raw&&!payload) return;
  const full=payload?(raw?payload+'User comment: '+raw:payload+'Please analyse these files.'):raw;

  addMessage('user', raw||(payload?'📎 (files attached)':''));
  msgIn.value=''; msgIn.style.height='auto';
  document.getElementById('char-info').textContent='';
  isThinking=true; setThinking(true); lastCount++;

  try{
    await fetch('/api/send',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({message:full})});
    startPoll();
  }catch{ setThinking(false); isThinking=false; addMessage('system','COMMS FAILURE — CHECK SERVER'); }
}

async function sendImagine(){
  if(isThinking) return;
  const raw=msgIn.value.trim();
  const payload=buildFilePayload();
  if(!raw&&!payload) return;
  const base=payload?(raw?payload+'User comment: '+raw:payload+'Analyse these.'):raw;

  addMessage('user','✨ [AESTHETIC] '+(raw||'(files)'));
  msgIn.value=''; msgIn.style.height='auto';
  isThinking=true; setThinking(true); lastCount++;

  try{
    await fetch('/api/imagine',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({message:base})});
    startPoll();
  }catch{ setThinking(false); isThinking=false; addMessage('system',' OFFLINE'); }
}

function useSuggestion(chip){ msgIn.value=chip.textContent.trim(); sendMessage(); }

async function clearChat(){
  if(!confirm('CLEAR ALL COMPUTER LOGS?')) return;
  await fetch('/api/clear',{method:'POST'});
  conv.innerHTML=
    '<div class="welcome"><div class="delta">🛡️</div>'+
    '<div class="w-title">COMPUTER ONLINE</div>'+
    '<div class="w-sub">LOGS CLEARED — AWAITING INQUIRY</div></div>';
  lastCount=0; lastContent=''; toast('✓ Logs cleared');
}

// ═══════════════════════════════════════════════════════════
//  LOAD HISTORY
// ═══════════════════════════════════════════════════════════
async function loadHistory(){
  try{
    const r=await fetch('/api/history'); const hist=await r.json();
    const newLast=hist.length?hist[hist.length-1].content:'';
    if((hist.length!==lastCount||newLast!==lastContent)&&!isThinking){
      lastContent=newLast;
      const prev=lastCount;
      const sp=conv.scrollHeight-conv.scrollTop;
      conv.innerHTML='';
      hist.forEach(m=>addMessage(m.role,m.content,false));
      lastCount=hist.length;
      if(!userScrolled) conv.scrollTop=conv.scrollHeight;
      else conv.scrollTop=conv.scrollHeight-sp;
      // TTS on new assistant message
      if(hist.length>prev){
        const last=hist[hist.length-1];
        if(last.role==='assistant') playTTS(last.content);
      }
    }
  }catch(e){}
}

// ═══════════════════════════════════════════════════════════
//  DRAG-AND-DROP UPLOAD
// ═══════════════════════════════════════════════════════════
const dropZone =document.getElementById('drop-zone');
const fileInput=document.getElementById('file-input');
const filePrev =document.getElementById('file-prev');

['dragenter','dragover','dragleave','drop'].forEach(ev=>{
  dropZone.addEventListener(ev,e=>{e.preventDefault();e.stopPropagation();});
  document.body.addEventListener(ev,e=>{e.preventDefault();e.stopPropagation();});
});
['dragenter','dragover'].forEach(ev=>dropZone.addEventListener(ev,()=>dropZone.classList.add('drag-over')));
['dragleave','drop'].forEach(ev=>dropZone.addEventListener(ev,()=>dropZone.classList.remove('drag-over')));
dropZone.addEventListener('drop',e=>uploadFiles(e.dataTransfer.files));
dropZone.addEventListener('click',()=>fileInput.click());
fileInput.addEventListener('change',e=>{ uploadFiles(e.target.files); fileInput.value=''; });

async function uploadFiles(files){
  if(!files.length) return;
  const fd=new FormData();
  for(const f of files) fd.append('files',f);
  toast('⏳ UPLOADING...');
  try{
    const r=await fetch('/api/upload',{method:'POST',body:fd});
    const d=await r.json();
    if(d.status==='success'){
      d.files.forEach(f=>pendingFiles.push(f));
      showFilePreviews(d.files);
      toast(`✓ ${d.files.length} file(s) staged — add comment and TRANSMIT`);
    }
  }catch{ toast('❌ Upload failed',true); }
}

function showFilePreviews(files){
  files.forEach((f,i)=>{
    const idx=pendingFiles.length-files.length+i;
    const icon=f.original_name.match(/\.pdf$/i)?'📄':
               f.original_name.match(/\.(png|jpg|jpeg|gif|webp)$/i)?'🖼':
               f.original_name.match(/\.xlsx?$/i)?'📊':
               f.original_name.match(/\.html?$/i)?'🌐':'📎';
    const el=document.createElement('div'); el.className='fpi';
    el.innerHTML=`<span>${icon}</span><span>${f.original_name} (${(f.size/1024).toFixed(1)}KB)</span>`+
      `<button onclick="removePending(${idx},this.parentElement)">✕</button>`;
    filePrev.appendChild(el);
  });
}
function removePending(idx,el){ pendingFiles[idx]=null; el.remove(); }

function buildFilePayload(){
  const active=pendingFiles.filter(f=>f!==null);
  if(!active.length) return null;
  let p=`I've attached ${active.length} file(s):\n\n`;
  active.forEach(f=>{
    p+=`**File: ${f.original_name}**\nSize: ${(f.size/1024).toFixed(1)} KB\nContent:\n\`\`\`\n${f.preview}\n\`\`\`\n\n`;
  });
  pendingFiles=[]; filePrev.innerHTML=''; return p;
}

// ═══════════════════════════════════════════════════════════
//  TEXTAREA EVENTS
// ═══════════════════════════════════════════════════════════
msgIn.addEventListener('input',function(){
  this.style.height='auto';
  this.style.height=Math.min(this.scrollHeight,130)+'px';
  const n=this.value.length;
  const ci=document.getElementById('char-info');
  ci.textContent=n>200?n+' / 50000':'';
  ci.style.color=n>45000?'var(--coral)':'var(--dimteal)';
});
msgIn.addEventListener('keydown',e=>{
  if(e.key==='Enter'&&!e.shiftKey){ e.preventDefault(); sendMessage(); }
});
conv.addEventListener('scroll',()=>{
  const atBot=conv.scrollHeight-conv.scrollTop-conv.clientHeight<60;
  userScrolled=!atBot; if(atBot) userScrolled=false;
});

// ═══════════════════════════════════════════════════════════
//  INIT
// ═══════════════════════════════════════════════════════════
loadHistory();
setInterval(loadHistory,1500);
setTimeout(()=>msgIn.focus(),400);
toast('🖖 NOVA COMPUTER ONLINE');
</script>
</body>
</html>
"""